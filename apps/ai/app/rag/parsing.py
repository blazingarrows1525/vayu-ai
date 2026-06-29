"""Extract plain text from uploaded documents."""

from __future__ import annotations

import io
from html.parser import HTMLParser

_EXT_TO_TYPE = {
    "pdf": "pdf",
    "docx": "docx",
    "pptx": "pptx",
    "md": "md",
    "markdown": "md",
    "txt": "txt",
    "csv": "csv",
    "xlsx": "xlsx",
}


def _ext(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def detect_source_type(filename: str) -> str:
    return _EXT_TO_TYPE.get(_ext(filename), "txt")


def _xlsx_to_text(data: bytes) -> str:
    """Flatten every sheet to `cell | cell | cell` rows so a spreadsheet is
    chunked + embedded like any other document."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        out: list[str] = []
        for ws in wb.worksheets:
            out.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append(" | ".join(cells))
        return "\n".join(out)
    finally:
        wb.close()


def _pptx_to_text(data: bytes) -> str:
    """Flatten a slide deck: each slide's shape text, table cells, and speaker
    notes become plain text under a `# Slide N` header so a deck is chunked +
    embedded like any other document."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        out.append(f"# Slide {index}")
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    out.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"Notes: {notes}")
    return "\n".join(out)


def extract_text(filename: str, data: bytes) -> str:
    ext = _ext(filename)
    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join((page.extract_text() or "") for page in reader.pages)
    if ext == "docx":
        import docx

        document = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in document.paragraphs)
    if ext == "pptx":
        return _pptx_to_text(data)
    if ext == "xlsx":
        return _xlsx_to_text(data)
    if ext in {"txt", "md", "markdown", "csv", ""}:
        return data.decode("utf-8", errors="replace")
    raise ValueError(f"Unsupported file type: .{ext}")


_BLOCK_TAGS = {"p", "div", "br", "li", "h1", "h2", "h3", "h4", "section", "article", "tr"}
_DROP_TAGS = {"script", "style", "noscript", "template"}


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: object) -> None:
        if tag in _DROP_TAGS:
            self._skip += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in _DROP_TAGS and self._skip:
            self._skip -= 1
        elif tag in _BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip and data.strip():
            self.parts.append(data)


def html_to_text(html: str) -> str:
    parser = _TextExtractor()
    parser.feed(html)
    lines = [ln.strip() for ln in "".join(parser.parts).splitlines()]
    return "\n".join(ln for ln in lines if ln)

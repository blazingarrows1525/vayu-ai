from fastapi.testclient import TestClient

from app.main import app
from app.rag.chunking import chunk_text
from app.rag.parsing import detect_source_type, extract_text

client = TestClient(app)


def test_rag_ask_requires_auth() -> None:
    assert client.post("/v1/rag/ask", json={"query": "hi"}).status_code == 401


def test_knowledge_list_requires_auth() -> None:
    assert client.get("/v1/knowledge").status_code == 401


def test_chunking_splits_long_text() -> None:
    text = "\n".join(f"Paragraph number {i} with some words." for i in range(200))
    chunks = chunk_text(text, target_chars=400, overlap=40)
    assert len(chunks) > 1
    assert all(c.strip() for c in chunks)


def test_parsing_plaintext_and_type() -> None:
    assert detect_source_type("notes.md") == "md"
    assert extract_text("notes.txt", b"hello world") == "hello world"


def test_pptx_extraction() -> None:
    import io

    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(0, 0, 914400, 914400)  # 1in x 1in (EMU)
    box.text_frame.text = "Quarterly revenue grew 30 percent"
    buf = io.BytesIO()
    prs.save(buf)

    assert detect_source_type("deck.pptx") == "pptx"
    text = extract_text("deck.pptx", buf.getvalue())
    assert "Quarterly revenue grew 30 percent" in text
